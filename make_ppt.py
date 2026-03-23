import json
import argparse
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


BG = RGBColor(0xF0, 0xF8, 0xFF)        # alice blue (near white)
ACCENT = RGBColor(0x26, 0x9B, 0xD4)    # sky blue
ACCENT2 = RGBColor(0x05, 0x6A, 0xA0)   # deeper blue for subtitle
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x44, 0x66, 0x88)      # dark slate for sub-bullets
GOLD = RGBColor(0x0A, 0x5C, 0x8A)      # deep blue replaces gold
HEADER_BG = RGBColor(0x26, 0x9B, 0xD4) # sky blue header bar
DIVIDER = RGBColor(0xB0, 0xD8, 0xF0)   # light blue divider
TEXT = RGBColor(0x1A, 0x2E, 0x44)      # dark navy for main text on light bg

W, H = Inches(13.33), Inches(7.5)
BASE_DIR = Path(__file__).resolve().parent
CONTENT_PATH = BASE_DIR / "content.json"


def load_content(path: Path):
    with path.open("r", encoding="utf-8") as f:
        return json.load(f)


def create_presentation():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    return prs


prs = create_presentation()
blank_layout = prs.slide_layouts[6]


def set_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, l, t, w, h, color):
    shape = slide.shapes.add_shape(1, l, t, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h, size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb


def title_slide(title, subtitle=""):
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide)
    add_rect(slide, Inches(0), Inches(0), Inches(0.18), H, ACCENT)
    add_text(slide, title, Inches(0.5), Inches(2.2), Inches(12), Inches(1.8), size=42, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
    if subtitle:
        add_text(slide, subtitle, Inches(0.5), Inches(4.1), Inches(12), Inches(0.8), size=24, color=ACCENT2, align=PP_ALIGN.CENTER)
    add_rect(slide, Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.06), ACCENT)


def section_slide(title, subtitle=""):
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide)
    add_rect(slide, Inches(0), Inches(0), W, Inches(0.12), ACCENT)
    add_rect(slide, Inches(0), H - Inches(0.12), W, Inches(0.12), ACCENT)
    add_text(slide, title, Inches(1), Inches(2.8), Inches(11), Inches(1.2), size=36, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    if subtitle:
        add_text(slide, subtitle, Inches(1), Inches(4.1), Inches(11), Inches(0.8), size=22, color=ACCENT2, align=PP_ALIGN.CENTER)


def content_slide(title, bullets):
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide)
    add_rect(slide, Inches(0), Inches(0), W, Inches(1.1), HEADER_BG)
    add_text(slide, title, Inches(0.3), Inches(0.15), Inches(12.5), Inches(0.8), size=28, bold=True, color=WHITE)
    add_rect(slide, Inches(0.3), Inches(1.15), Inches(12.5), Inches(0.04), ACCENT)

    txb = slide.shapes.add_textbox(Inches(0.4), Inches(1.35), Inches(12.4), Inches(5.8))
    tf = txb.text_frame
    tf.word_wrap = True

    for index, bullet in enumerate(bullets):
        level, text = bullet
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.level = level
        p.space_before = Pt(4 if level == 0 else 2)
        indent = "    " * level
        bullet_char = "▸ " if level == 0 else "• "
        run = p.add_run()
        run.text = indent + bullet_char + text
        run.font.size = Pt(18 if level == 0 else 16)
        run.font.bold = level == 0
        run.font.color.rgb = TEXT if level == 0 else GRAY


def two_col_slide(title, left_items, right_items, left_head="", right_head=""):
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide)
    add_rect(slide, Inches(0), Inches(0), W, Inches(1.1), HEADER_BG)
    add_text(slide, title, Inches(0.3), Inches(0.15), Inches(12.5), Inches(0.8), size=28, bold=True, color=WHITE)
    add_rect(slide, Inches(0.3), Inches(1.15), Inches(12.5), Inches(0.04), ACCENT)
    add_rect(slide, Inches(6.6), Inches(1.25), Inches(0.04), Inches(5.9), DIVIDER)

    columns = [
        (left_head, left_items, Inches(0.4)),
        (right_head, right_items, Inches(6.8)),
    ]

    for head, items, lx in columns:
        cw = Inches(6.0)
        if head:
            add_text(slide, head, lx, Inches(1.25), cw, Inches(0.4), size=18, bold=True, color=ACCENT2)
        txb = slide.shapes.add_textbox(lx, Inches(1.7), cw, Inches(5.5))
        tf = txb.text_frame
        tf.word_wrap = True
        for index, item in enumerate(items):
            p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
            p.space_before = Pt(5)
            run = p.add_run()
            run.text = "▸ " + item
            run.font.size = Pt(16)
            run.font.color.rgb = TEXT


# ── Additional slide type variants ─────────────────────────────

def title_slide_centered(title, subtitle=""):
    """Full-screen centered title with large accent circle behind text."""
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, BG)
    # large decorative circle
    from pptx.enum.shapes import MSO_SHAPE
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(4.2), Inches(1.0), Inches(5), Inches(5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(0xDB, 0xEE, 0xF9)
    circle.line.fill.background()
    add_text(slide, title, Inches(1), Inches(2.5), Inches(11.3), Inches(1.5),
             size=44, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
    if subtitle:
        add_text(slide, subtitle, Inches(1), Inches(4.2), Inches(11.3), Inches(1.0),
                 size=22, color=ACCENT2, align=PP_ALIGN.CENTER)


def title_slide_split(title, subtitle=""):
    """Left half accent-colored, right half white. Title on left, subtitle on right."""
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, WHITE)
    add_rect(slide, Inches(0), Inches(0), Inches(6.66), H, ACCENT)
    add_text(slide, title, Inches(0.5), Inches(2.5), Inches(5.8), Inches(2.0),
             size=40, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle, Inches(7.2), Inches(2.8), Inches(5.5), Inches(1.5),
                 size=22, color=GRAY, align=PP_ALIGN.LEFT)
    add_rect(slide, Inches(6.66), Inches(2.0), Inches(0.06), Inches(3.5), ACCENT)


def section_slide_banner(title, subtitle=""):
    """Full-width accent banner in the middle with white text."""
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, WHITE)
    add_rect(slide, Inches(0), Inches(2.5), W, Inches(2.5), ACCENT)
    add_text(slide, title, Inches(1), Inches(2.7), Inches(11.3), Inches(1.2),
             size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if subtitle:
        add_text(slide, subtitle, Inches(1), Inches(3.9), Inches(11.3), Inches(0.8),
                 size=20, color=RGBColor(0xDB, 0xEE, 0xF9), align=PP_ALIGN.CENTER)


def section_slide_corner(title, subtitle=""):
    """Accent corner triangle decoration with left-aligned title."""
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, BG)
    # top-left accent block
    add_rect(slide, Inches(0), Inches(0), Inches(0.5), Inches(0.5), ACCENT)
    add_rect(slide, Inches(0), Inches(0), Inches(1.5), Inches(0.15), ACCENT)
    add_rect(slide, Inches(0), Inches(0), Inches(0.15), Inches(1.5), ACCENT)
    # bottom-right accent block
    add_rect(slide, W - Inches(0.5), H - Inches(0.5), Inches(0.5), Inches(0.5), ACCENT)
    add_rect(slide, W - Inches(1.5), H - Inches(0.15), Inches(1.5), Inches(0.15), ACCENT)
    add_rect(slide, W - Inches(0.15), H - Inches(1.5), Inches(0.15), Inches(1.5), ACCENT)
    add_text(slide, title, Inches(1.5), Inches(2.8), Inches(10), Inches(1.2),
             size=36, bold=True, color=TEXT, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle, Inches(1.5), Inches(4.1), Inches(10), Inches(0.8),
                 size=20, color=GRAY, align=PP_ALIGN.LEFT)


def content_slide_cards(title, bullets):
    """Each level-0 bullet becomes a card-like box with its sub-bullets inside."""
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, WHITE)
    add_text(slide, title, Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.8),
             size=28, bold=True, color=ACCENT)
    add_rect(slide, Inches(0.4), Inches(1.05), Inches(12.5), Inches(0.04), ACCENT)

    # group bullets into cards: each level-0 starts a new card
    cards = []
    for level, text in bullets:
        if level == 0:
            cards.append({"title": text, "items": []})
        else:
            if cards:
                cards[-1]["items"].append(text)

    cols = 2
    card_w = Inches(6.0)
    card_h = Inches(2.8)
    x_starts = [Inches(0.4), Inches(6.9)]
    y_start = Inches(1.3)
    y_gap = Inches(0.15)

    for i, card in enumerate(cards):
        col = i % cols
        row = i // cols
        x = x_starts[col]
        y = y_start + row * (card_h + y_gap)
        if y + card_h > H:
            break
        # card background
        add_rect(slide, x, y, card_w, card_h, RGBColor(0xE8, 0xF4, 0xFB))
        # card accent bar
        add_rect(slide, x, y, Inches(0.08), card_h, ACCENT)
        # card title
        add_text(slide, "▸ " + card["title"], x + Inches(0.2), y + Inches(0.1),
                 card_w - Inches(0.3), Inches(0.5), size=16, bold=True, color=TEXT)
        # card items
        if card["items"]:
            txb = slide.shapes.add_textbox(x + Inches(0.3), y + Inches(0.6),
                                           card_w - Inches(0.5), card_h - Inches(0.7))
            tf = txb.text_frame
            tf.word_wrap = True
            for j, item in enumerate(card["items"]):
                p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                p.space_before = Pt(2)
                run = p.add_run()
                run.text = "• " + item
                run.font.size = Pt(13)
                run.font.color.rgb = GRAY


def content_slide_numbered(title, bullets):
    """Level-0 bullets get large numbered circles; sub-bullets indented below."""
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, BG)
    add_text(slide, title, Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.8),
             size=28, bold=True, color=ACCENT)
    add_rect(slide, Inches(0.4), Inches(1.05), Inches(12.5), Inches(0.04), ACCENT)

    y = Inches(1.3)
    num = 0
    for level, text in bullets:
        if y > H - Inches(0.5):
            break
        if level == 0:
            num += 1
            from pptx.enum.shapes import MSO_SHAPE
            circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.4), y, Inches(0.5), Inches(0.5))
            circle.fill.solid()
            circle.fill.fore_color.rgb = ACCENT
            circle.line.fill.background()
            # number inside circle
            tf = circle.text_frame
            tf.word_wrap = False
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = str(num)
            run.font.size = Pt(16)
            run.font.bold = True
            run.font.color.rgb = WHITE
            add_text(slide, text, Inches(1.1), y + Inches(0.05), Inches(11.5), Inches(0.45),
                     size=18, bold=True, color=TEXT)
            y += Inches(0.55)
        else:
            add_text(slide, "• " + text, Inches(1.3), y, Inches(11.3), Inches(0.4),
                     size=15, color=GRAY)
            y += Inches(0.4)


def content_slide_sidebar(title, bullets):
    """Left accent sidebar with title rotated-style, bullets on the right."""
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, WHITE)
    # sidebar
    add_rect(slide, Inches(0), Inches(0), Inches(2.5), H, ACCENT)
    add_text(slide, title, Inches(0.15), Inches(1.5), Inches(2.2), Inches(4.5),
             size=22, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    txb = slide.shapes.add_textbox(Inches(2.8), Inches(0.5), Inches(10.0), Inches(6.5))
    tf = txb.text_frame
    tf.word_wrap = True

    for index, (level, text) in enumerate(bullets):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.level = level
        p.space_before = Pt(6 if level == 0 else 3)
        indent = "    " * level
        bullet_char = "▸ " if level == 0 else "• "
        run = p.add_run()
        run.text = indent + bullet_char + text
        run.font.size = Pt(17 if level == 0 else 15)
        run.font.bold = level == 0
        run.font.color.rgb = TEXT if level == 0 else GRAY


def two_col_slide_boxed(title, left_items, right_items, left_head="", right_head=""):
    """Two columns with colored card backgrounds."""
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, WHITE)
    add_text(slide, title, Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.8),
             size=28, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_rect(slide, Inches(0.4), Inches(1.05), Inches(12.5), Inches(0.04), ACCENT)

    col_specs = [
        (left_head, left_items, Inches(0.4), RGBColor(0xE8, 0xF4, 0xFB)),
        (right_head, right_items, Inches(6.9), RGBColor(0xF0, 0xF0, 0xF8)),
    ]
    for head, items, lx, bg_color in col_specs:
        add_rect(slide, lx, Inches(1.2), Inches(6.0), Inches(5.9), bg_color)
        add_rect(slide, lx, Inches(1.2), Inches(6.0), Inches(0.06), ACCENT)
        if head:
            add_text(slide, head, lx + Inches(0.2), Inches(1.35), Inches(5.5), Inches(0.4),
                     size=17, bold=True, color=ACCENT2)
        txb = slide.shapes.add_textbox(lx + Inches(0.2), Inches(1.85), Inches(5.5), Inches(5.0))
        tf = txb.text_frame
        tf.word_wrap = True
        for index, item in enumerate(items):
            p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
            p.space_before = Pt(6)
            run = p.add_run()
            run.text = "▸ " + item
            run.font.size = Pt(15)
            run.font.color.rgb = TEXT


def quote_slide(text, attribution=""):
    """Large centered quote with optional attribution."""
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, BG)
    add_rect(slide, Inches(2), Inches(1.5), Inches(0.08), Inches(4.5), ACCENT)
    add_text(slide, '"', Inches(1.5), Inches(1.2), Inches(1), Inches(1),
             size=72, bold=True, color=RGBColor(0xDB, 0xEE, 0xF9))
    add_text(slide, text, Inches(2.5), Inches(2.0), Inches(9), Inches(3.0),
             size=24, bold=False, color=TEXT, align=PP_ALIGN.LEFT)
    if attribution:
        add_text(slide, "— " + attribution, Inches(2.5), Inches(5.2), Inches(9), Inches(0.5),
                 size=18, bold=True, color=ACCENT2, align=PP_ALIGN.LEFT)


def three_col_slide(title, col1, col2, col3, head1="", head2="", head3=""):
    """Three equal columns."""
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, BG)
    add_text(slide, title, Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.8),
             size=28, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_rect(slide, Inches(0.4), Inches(1.05), Inches(12.5), Inches(0.04), ACCENT)

    col_w = Inches(3.8)
    x_starts = [Inches(0.4), Inches(4.6), Inches(8.8)]
    cols_data = [(head1, col1), (head2, col2), (head3, col3)]

    for x, (head, items) in zip(x_starts, cols_data):
        add_rect(slide, x, Inches(1.2), col_w, Inches(0.06), ACCENT)
        if head:
            add_text(slide, head, x, Inches(1.35), col_w, Inches(0.4),
                     size=16, bold=True, color=ACCENT2, align=PP_ALIGN.CENTER)
        txb = slide.shapes.add_textbox(x + Inches(0.1), Inches(1.85), col_w - Inches(0.2), Inches(5.2))
        tf = txb.text_frame
        tf.word_wrap = True
        for index, item in enumerate(items):
            p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
            p.space_before = Pt(5)
            run = p.add_run()
            run.text = "▸ " + item
            run.font.size = Pt(14)
            run.font.color.rgb = TEXT


def big_number_slide(number, label, description=""):
    """Hero slide with a large number/stat and label."""
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, ACCENT)
    add_text(slide, str(number), Inches(1), Inches(1.5), Inches(11.3), Inches(2.5),
             size=96, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, label, Inches(1), Inches(4.0), Inches(11.3), Inches(1.0),
             size=32, bold=True, color=RGBColor(0xDB, 0xEE, 0xF9), align=PP_ALIGN.CENTER)
    if description:
        add_text(slide, description, Inches(2), Inches(5.2), Inches(9.3), Inches(1.0),
                 size=20, color=WHITE, align=PP_ALIGN.CENTER)


def build_slides(spec):
    for slide in spec["slides"]:
        slide_type = slide["type"]
        if slide_type == "title":
            title_slide(slide["title"], slide.get("subtitle", ""))
        elif slide_type == "title_centered":
            title_slide_centered(slide["title"], slide.get("subtitle", ""))
        elif slide_type == "title_split":
            title_slide_split(slide["title"], slide.get("subtitle", ""))
        elif slide_type == "section":
            section_slide(slide["title"], slide.get("subtitle", ""))
        elif slide_type == "section_banner":
            section_slide_banner(slide["title"], slide.get("subtitle", ""))
        elif slide_type == "section_corner":
            section_slide_corner(slide["title"], slide.get("subtitle", ""))
        elif slide_type == "content":
            content_slide(slide["title"], slide["bullets"])
        elif slide_type == "content_cards":
            content_slide_cards(slide["title"], slide["bullets"])
        elif slide_type == "content_numbered":
            content_slide_numbered(slide["title"], slide["bullets"])
        elif slide_type == "content_sidebar":
            content_slide_sidebar(slide["title"], slide["bullets"])
        elif slide_type == "two_col":
            two_col_slide(
                slide["title"],
                slide["left"],
                slide["right"],
                slide.get("left_head", ""),
                slide.get("right_head", ""),
            )
        elif slide_type == "two_col_boxed":
            two_col_slide_boxed(
                slide["title"],
                slide["left"],
                slide["right"],
                slide.get("left_head", ""),
                slide.get("right_head", ""),
            )
        elif slide_type == "quote":
            quote_slide(slide["text"], slide.get("attribution", ""))
        elif slide_type == "three_col":
            three_col_slide(
                slide["title"],
                slide["col1"],
                slide["col2"],
                slide["col3"],
                slide.get("head1", ""),
                slide.get("head2", ""),
                slide.get("head3", ""),
            )
        elif slide_type == "big_number":
            big_number_slide(slide["number"], slide["label"], slide.get("description", ""))
        else:
            raise ValueError(f"Unsupported slide type: {slide_type}")


def main():
    parser = argparse.ArgumentParser(description="Generate a PPTX from a content JSON file.")
    parser.add_argument("-i", "--input", default=str(CONTENT_PATH),
                        help="Path to content JSON file (default: content.json)")
    parser.add_argument("-o", "--output", default=None,
                        help="Path for output PPTX (default: from JSON or protein_design_survey.pptx)")
    args = parser.parse_args()

    spec = load_content(Path(args.input))
    build_slides(spec)

    if args.output:
        output_path = Path(args.output)
    else:
        try:
            output_path = Path(spec.get("output", BASE_DIR / "protein_design_survey.pptx"))
        except Exception as e:
            print(f"Error occurred while setting output path: {e}")
            return

    prs.save(str(output_path))
    print(f"Saved: {output_path} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
